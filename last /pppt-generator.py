import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import OpenAI
import requests
import os

# ========== SET API KEY ==========
OPENAI_API_KEY = "YOUR_API_KEY"
client = OpenAI(api_key=OPENAI_API_KEY)

# ========== AI SLIDE GENERATOR ==========
def generate_slides(topic, num_slides):

    prompt = f"""
    Create a professional presentation on: {topic}

    Give exactly {num_slides} slides.

    For each slide provide:
    Slide Title:
    Bullet Points (3-5)
    Speaker Notes

    Format strictly like:

    Slide 1:
    Title:
    Points:
    Notes:

    Slide 2:
    ...
    """
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role":"user","content":prompt}]
    )

    return response.choices[0].message.content


# ========== PPT BUILDER ==========
def build_ppt(content):

    prs = Presentation()

    slides = content.split("Slide ")

    for slide in slides[1:]:

        lines = slide.split("\n")
        title = ""
        points = []
        notes = ""

        mode = ""
        for line in lines:
            if "Title:" in line:
                title = line.replace("Title:","").strip()
            elif "Points:" in line:
                mode = "points"
            elif "Notes:" in line:
                mode = "notes"
            else:
                if mode == "points":
                    points.append(line.strip())
                elif mode == "notes":
                    notes += line.strip() + " "

        slide_layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)

        slide_obj.shapes.title.text = title
        body = slide_obj.shapes.placeholders[1]
        tf = body.text_frame

        for p in points:
            if p:
                tf.add_paragraph().text = p

        slide_obj.notes_slide.notes_text_frame.text = notes

    file_path = "generated_presentation.pptx"
    prs.save(file_path)
    return file_path
