import streamlit as st
from pptx import Presentation
from openai import OpenAI
from pypdf import PdfReader

OPENAI_API_KEY = "YOUR_API_KEY"
client = OpenAI(api_key=OPENAI_API_KEY)

# -------- Extract Text from PDF --------
def extract_text(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# -------- AI Structuring --------
def generate_slides_from_doc(text, num_slides):

    prompt = f"""
    Convert this document into a {num_slides}-slide presentation.

    Extract key ideas and structure logically.

    Format:

    Slide 1:
    Title:
    Points:
    Notes:

    Slide 2:
    ...
    """

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role":"user","content":prompt + text[:6000]}]
    )

    return response.choices[0].message.content
# -------- PPT Builder --------
def build_ppt(content):

    prs = Presentation()
    slides = content.split("Slide ")

    for slide in slides[1:]:

        lines = slide.split("\n")
        title = ""
        points = []
        notes = ""
        mode = ""

        for line in lines:
            if "Title:" in line:
                title = line.replace("Title:","").strip()
            elif "Points:" in line:
                mode = "points"
            elif "Notes:" in line:
                mode = "notes"
            else:
                if mode == "points":
                    points.append(line.strip())
                elif mode == "notes":
                    notes += line.strip() + " "

        slide_layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(slide_layout)
